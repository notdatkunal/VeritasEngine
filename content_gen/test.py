from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


SLIDES_DATA = [
    {
        "slide_number": 1,
        "slide_type": "Hook",
        "headline": "AI is taking your job? Think again. 🤯",
        "body_text": "Swipe to see what *really* changed in tech."
    },
    {
        "slide_number": 2,
        "slide_type": "Second Chance Cover",
        "headline": "Netflix just proved it: Efficiency > Genius. 💸",
        "body_text": "Major tech players are ditching tools not because they're bad, but because they cost too much or move too slowly. This isn't about *if* AI is powerful—it’s about the **economic reality** of building software today."
    },
    {
        "slide_number": 3,
        "slide_type": "Correction",
        "headline": "Stop Reading Between the Lines! 🛑",
        "body_text": "AI isn't replacing developers. It’s changing what \"developer\" means. The skill shifts from **manual coding** to **strategic direction.**"
    },
    {
        "slide_number": 4,
        "slide_type": "Thesis",
        "headline": "The New Currency of Engineering 💰",
        "body_text": "Value isn't in lines of code. It’s in **Engineering Hours Saved.** AI tools maximize output with minimal manual effort."
    },
    {
        "slide_number": 5,
        "slide_type": "Framework (Savable Value)",
        "headline": "🛠️ The AI Efficiency Audit (Save This!)",
        "body_text": "1. **Identify Repetition:** What do you code manually every week?\n2. **Measure Time:** How many hours are spent on that task?\n3. **Automate/Delegate:** Can AI handle 80% of the effort?\n*That's your ROI.*"
    },
    {
        "slide_number": 6,
        "slide_type": "Historical Parallel",
        "headline": "It’s Not New. It's Just Faster. 🕰️",
        "body_text": "Think frameworks replacing manual HTML/CSS. Or Python simplifying C++. Tech always reduces *manual effort*, enabling bigger leaps."
    },
    {
        "slide_number": 7,
        "slide_type": "Mindset Shift",
        "headline": "Your Role is Now an Architect. 🏗️",
        "body_text": "Don't be a coder. Be the **System Designer.** Your job: Define the problem and structure the solution. AI handles the plumbing; you design the skyscraper."
    },
    {
        "slide_number": 8,
        "slide_type": "Call-to-Action",
        "headline": "Don't Just Swipe. ACT. 👇",
        "body_text": "💾 **SAVE THIS POST:** Use the \"AI Efficiency Audit\" framework on your next project plan.\n\n📤 **SHARE TO DM:** Send this to a skeptical manager or junior developer who needs this perspective!"
    }
]


def create_presentation(file_name: str = "Instagram_Carousel.pptx") -> None:
    """
    Generates a beautifully formatted square PowerPoint presentation
    optimized for an Instagram Carousel.
    """
    prs = Presentation()
    
    # Set slide dimensions to 10x10 inches (1:1 aspect ratio for Instagram)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(10)

    # Theme colors
    bg_color = RGBColor(248, 249, 250)    # Light gray/off-white background
    title_color = RGBColor(33, 37, 41)    # Dark text for contrast
    body_color = RGBColor(73, 80, 87)     # Slightly softer dark gray for body

    for slide_data in SLIDES_DATA:
        # Layout 6 is a blank slide; avoiding placeholders that clip on resized canvases
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # Format Headline
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(2.0))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.text = slide_data["headline"]
        
        for paragraph in title_tf.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.bold = True
            paragraph.font.size = Pt(44)
            paragraph.font.color.rgb = title_color

        # Format Body Text
        body_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(8.0), Inches(5.0))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        body_tf.text = slide_data["body_text"]
        
        for paragraph in body_tf.paragraphs:
            paragraph.font.size = Pt(28)
            paragraph.font.color.rgb = body_color
            # Optional: adjust alignment if you prefer center-aligned body text
            # paragraph.alignment = PP_ALIGN.CENTER

    # Save to the current directory
    prs.save(file_name)
    print(f"Success: Presentation saved as {file_name}")


if __name__ == "__main__":
    create_presentation()